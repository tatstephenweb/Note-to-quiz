from flask import Flask, render_template, request, session, redirect, url_for
import json, os, fitz
from pptx import Presentation
from groq import Groq
from dotenv import load_dotenv

app = Flask(__name__)
app.secret_key = "quiz_secret_key"
UPLOAD_FOLDER = "uploads"

load_dotenv()
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

client = Groq(api_key=os.getenv("GROQ_API_KEY"))

# ── Extract text from uploaded file ──────────────────────────
def extract_text(filepath):
    if filepath.endswith(".pdf"):
        doc = fitz.open(filepath)
        return "\n".join(page.get_text() for page in doc)

    elif filepath.endswith(".pptx"):
        prs = Presentation(filepath)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text.append(shape.text_frame.text)
        return "\n".join(text)

    return ""


def generate_questions(text):
    prompt = f"""
Generate 10 multiple choice questions from this lecture content.
Return ONLY a valid JSON array, no explanation, no markdown, no backticks.
Each object must have exactly these fields:
- "id": number
- "question": string
- "options": array of 4 strings
- "answer": string (must exactly match one of the options)

Lecture content:
{text[:8000]}
"""
    response = client.chat.completions.create(
        model="llama-3.3-70b-versatile",
        messages=[{"role": "user", "content": prompt}]
    )
    raw = response.choices[0].message.content.strip()

    if raw.startswith("```"):
        raw = raw.split("\n", 1)[1]
    if raw.endswith("```"):
        raw = raw.rsplit("```", 1)[0]

    return json.loads(raw.strip())

# ── Routes ────────────────────────────────────────────────────
@app.route("/", methods=["GET", "POST"])
def upload():
    if request.method == "POST":
        file = request.files.get("file")

        if not file or file.filename == "":
            return render_template("upload.html", error="Please select a file.")

        filename = file.filename
        if not (filename.endswith(".pdf") or filename.endswith(".pptx")):
            return render_template("upload.html", error="Only PDF and PPTX files are supported.")

        filepath = os.path.join(UPLOAD_FOLDER, filename)
        file.save(filepath)

        try:
            text = extract_text(filepath)
            if not text.strip():
                return render_template("upload.html", error="Could not extract text from file. Try another file.")

            questions = generate_questions(text)

            with open("questions.json", "w") as f:
                json.dump(questions, f, indent=2)
            
            #session["questions"] = questions

            return redirect(url_for("start"))

        except Exception as e:
            return render_template("upload.html", error=f"Something went wrong: {str(e)}")

    return render_template("upload.html")

@app.route("/generated")
def start():
    session["current"] = 0
    session["score"] = 0
    session["answers"] = []
    return render_template("quiz-generated.html")

@app.route("/quiz", methods=["GET", "POST"])
def quiz():
    questions = load_questions()
    index = session.get("current", 0)

    if request.method == "POST":
        selected = request.form.get("answer")
        correct = questions[index]["answer"]

        session["answers"] = session.get("answers", []) + [selected]

        if selected == correct:
            session["score"] += 1

        session["current"] += 1
        index = session["current"]

        if index >= len(questions):
            return redirect(url_for("result"))

    question = questions[index]
    total = len(questions)
    return render_template("quiz.html", question=question, index=index, total=total)

@app.route("/result")
def result():
    score = session.get("score", 0)
    total = len(load_questions())
    return render_template("result.html", score=score, total=total)

@app.route("/home")
def home():
    return render_template("upload.html")

def load_questions():
    with open("questions.json") as f:
        return json.load(f)

if __name__ == "__main__":
    app.run(debug=True)